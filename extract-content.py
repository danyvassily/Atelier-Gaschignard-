#!/usr/bin/env python3
"""
Script temporaire pour extraire le texte des documents PDF et PPTX
"""

import sys
import os

def extract_pdf_text(pdf_path):
    """Extrait le texte d'un PDF"""
    try:
        import PyPDF2
        text = ""
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text
    except ImportError:
        return "PyPDF2 non installé. Installer avec: pip install PyPDF2"
    except Exception as e:
        return f"Erreur lors de l'extraction: {str(e)}"

def extract_pptx_text(pptx_path):
    """Extrait le texte d'un PPTX"""
    try:
        from pptx import Presentation
        text = ""
        prs = Presentation(pptx_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except ImportError:
        return "python-pptx non installé. Installer avec: pip install python-pptx"
    except Exception as e:
        return f"Erreur lors de l'extraction: {str(e)}"

if __name__ == "__main__":
    pdf_path = "../Présentation Canva.pdf"
    pptx_path = "../Texte.pptx"
    
    if os.path.exists(pdf_path):
        print("=" * 60)
        print("CONTENU DU PDF - Présentation Canva")
        print("=" * 60)
        print(extract_pdf_text(pdf_path))
        print("\n")
    
    if os.path.exists(pptx_path):
        print("=" * 60)
        print("CONTENU DU PPTX - Texte")
        print("=" * 60)
        print(extract_pptx_text(pptx_path))

